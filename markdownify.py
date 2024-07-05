import streamlit as st
from pptx import Presentation
from markdownify import markdownify as md

def convert_pptx_to_md(pptx_file):
    prs = Presentation(pptx_file)
    md_content = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                md_content += md(shape.text) + "\n\n"

    return md_content

def main():
    st.title("PowerPoint to Markdown Converter")

    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
    if uploaded_file is not None:
        md_content = convert_pptx_to_md(uploaded_file)

        st.markdown("### Markdown Content")
        st.text_area("Markdown Output", md_content, height=300)

        st.download_button(
            label="Download Markdown",
            data=md_content,
            file_name="presentation.md",
            mime="text/markdown",
        )

if __name__ == "__main__":
    main()
